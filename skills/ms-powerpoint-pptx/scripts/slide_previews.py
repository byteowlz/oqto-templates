#!/usr/bin/env python3
"""
Generate visual thumbnail grids from PowerPoint presentations.

Produces grid layouts of slide previews with adjustable columns (maximum 6).
Each grid holds up to cols×(cols+1) images. For larger decks, multiple
numbered grid files are generated automatically.

Program outputs the names of all created files.

Output naming:
- Single grid: {prefix}.jpg (when slides fit in one grid)
- Multiple grids: {prefix}-1.jpg, {prefix}-2.jpg, etc.

Grid capacity by column count:
- 3 columns: max 12 slides (3×4)
- 4 columns: max 20 slides (4×5)
- 5 columns: max 30 slides (5×6) [default]
- 6 columns: max 42 slides (6×7)

Usage:
    python slide_previews.py input.pptx [output_prefix] [--cols N] [--outline-placeholders]

Examples:
    python slide_previews.py deck.pptx
    # Creates: previews.jpg (default prefix)
    # Outputs:
    #   Created 1 grid(s):
    #     - previews.jpg

    python slide_previews.py big-deck.pptx visual --cols 4
    # Creates: visual-1.jpg, visual-2.jpg, visual-3.jpg
    # Outputs:
    #   Created 3 grid(s):
    #     - visual-1.jpg
    #     - visual-2.jpg
    #     - visual-3.jpg

    python slide_previews.py template.pptx analysis --outline-placeholders
    # Creates grids with red outlines around text placeholders
"""

import argparse
import subprocess
import sys
import tempfile
from pathlib import Path

from inventory import extract_text_inventory
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

# Configuration
PREVIEW_WIDTH = 300  # Fixed preview width in pixels
EXPORT_DPI = 100  # DPI for PDF export
MAX_COLUMNS = 6  # Upper limit for columns
DEFAULT_COLUMNS = 5  # Standard column count
IMAGE_QUALITY = 95  # JPEG compression level

# Layout settings
SPACING = 20  # Gap between previews
BORDER_THICKNESS = 2  # Border around previews
TEXT_SCALE = 0.12  # Text size relative to preview width
TEXT_MARGIN_SCALE = 0.4  # Text margin relative to text size


def run():
    parser = argparse.ArgumentParser(
        description="Generate preview grids from PowerPoint slides."
    )
    parser.add_argument("source", help="Input PowerPoint file (.pptx)")
    parser.add_argument(
        "prefix",
        nargs="?",
        default="previews",
        help="Output prefix for images (default: previews, creates prefix.jpg or prefix-N.jpg)",
    )
    parser.add_argument(
        "--cols",
        type=int,
        default=DEFAULT_COLUMNS,
        help=f"Column count (default: {DEFAULT_COLUMNS}, max: {MAX_COLUMNS})",
    )
    parser.add_argument(
        "--outline-placeholders",
        action="store_true",
        help="Draw borders around text placeholders",
    )

    args = parser.parse_args()

    # Cap columns at maximum
    columns = min(args.cols, MAX_COLUMNS)
    if args.cols > MAX_COLUMNS:
        print(f"Note: Columns capped at {MAX_COLUMNS} (requested {args.cols})")

    # Verify input
    source_path = Path(args.source)
    if not source_path.exists() or source_path.suffix.lower() != ".pptx":
        print(f"Error: Invalid PowerPoint file: {args.source}")
        sys.exit(1)

    # Set output path (always JPG)
    target_path = Path(f"{args.prefix}.jpg")

    print(f"Processing: {args.source}")

    try:
        with tempfile.TemporaryDirectory() as tmp:
            # Capture text region coordinates if outlining enabled
            text_regions = None
            slide_size = None
            if args.outline_placeholders:
                print("Locating text regions...")
                text_regions, slide_size = fetch_text_regions(source_path)
                if text_regions:
                    print(f"Found text on {len(text_regions)} slides")

            # Export slides to images
            slide_files = export_to_images(source_path, Path(tmp), EXPORT_DPI)
            if not slide_files:
                print("Error: No slides found")
                sys.exit(1)

            print(f"Found {len(slide_files)} slides")

            # Build grids (max cols×(cols+1) images each)
            output_files = build_grid_collection(
                slide_files,
                columns,
                PREVIEW_WIDTH,
                target_path,
                text_regions,
                slide_size,
            )

            # Display results
            print(f"Created {len(output_files)} grid(s):")
            for file in output_files:
                print(f"  - {file}")

    except Exception as ex:
        print(f"Error: {ex}")
        sys.exit(1)


def make_hidden_slide_image(dimensions):
    """Generate placeholder image for hidden slides."""
    img = Image.new("RGB", dimensions, color="#F0F0F0")
    canvas = ImageDraw.Draw(img)
    stroke = max(5, min(dimensions) // 100)
    canvas.line([(0, 0), dimensions], fill="#CCCCCC", width=stroke)
    canvas.line([(dimensions[0], 0), (0, dimensions[1])], fill="#CCCCCC", width=stroke)
    return img


def fetch_text_regions(pptx_file):
    """Extract ALL text regions from the presentation.

    Returns (text_regions, slide_dimensions).
    text_regions maps slide indices to lists of regions.
    Each region has 'left', 'top', 'width', 'height' in inches.
    slide_dimensions is (width_inches, height_inches).
    """
    deck = Presentation(str(pptx_file))
    inventory = extract_text_inventory(pptx_file, deck)
    text_regions = {}

    # Convert EMU to inches
    width_inches = (deck.slide_width or 9144000) / 914400.0
    height_inches = (deck.slide_height or 5143500) / 914400.0

    for slide_key, shapes in inventory.items():
        slide_idx = int(slide_key.split("-")[1])
        regions = []

        for shape_key, shape_data in shapes.items():
            regions.append(
                {
                    "left": shape_data.left,
                    "top": shape_data.top,
                    "width": shape_data.width,
                    "height": shape_data.height,
                }
            )

        if regions:
            text_regions[slide_idx] = regions

    return text_regions, (width_inches, height_inches)


def export_to_images(pptx_file, tmp_folder, dpi):
    """Export PowerPoint to images via PDF, including hidden slides."""
    print("Analyzing presentation structure...")
    deck = Presentation(str(pptx_file))
    total_count = len(deck.slides)

    # Identify hidden slides (1-based)
    hidden = {
        idx + 1
        for idx, slide in enumerate(deck.slides)
        if slide.element.get("show") == "0"
    }

    print(f"Total slides: {total_count}")
    if hidden:
        print(f"Hidden slides: {sorted(hidden)}")

    pdf_file = tmp_folder / f"{pptx_file.stem}.pdf"

    # Export to PDF
    print("Exporting to PDF...")
    result = subprocess.run(
        [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(tmp_folder),
            str(pptx_file),
        ],
        capture_output=True,
        text=True,
    )
    if result.returncode != 0 or not pdf_file.exists():
        raise RuntimeError("PDF export failed")

    # Convert PDF to images
    print(f"Generating images at {dpi} DPI...")
    result = subprocess.run(
        ["pdftoppm", "-jpeg", "-r", str(dpi), str(pdf_file), str(tmp_folder / "slide")],
        capture_output=True,
        text=True,
    )
    if result.returncode != 0:
        raise RuntimeError("Image generation failed")

    visible_files = sorted(tmp_folder.glob("slide-*.jpg"))

    # Build complete list with placeholders for hidden slides
    all_files = []
    visible_idx = 0

    # Get dimensions from first visible slide
    if visible_files:
        with Image.open(visible_files[0]) as img:
            placeholder_dims = img.size
    else:
        placeholder_dims = (1920, 1080)

    for slide_num in range(1, total_count + 1):
        if slide_num in hidden:
            # Create placeholder for hidden slide
            placeholder_file = tmp_folder / f"hidden-{slide_num:03d}.jpg"
            placeholder_img = make_hidden_slide_image(placeholder_dims)
            placeholder_img.save(placeholder_file, "JPEG")
            all_files.append(placeholder_file)
        else:
            # Use actual visible slide
            if visible_idx < len(visible_files):
                all_files.append(visible_files[visible_idx])
                visible_idx += 1

    return all_files


def build_grid_collection(
    image_files,
    columns,
    width,
    target_path,
    text_regions=None,
    slide_size=None,
):
    """Generate multiple preview grids, max cols×(cols+1) images each."""
    max_per_grid = columns * (columns + 1)
    output_files = []

    print(f"Building grids with {columns} columns (max {max_per_grid} images per grid)")

    # Process in chunks
    for chunk_idx, start in enumerate(range(0, len(image_files), max_per_grid)):
        end = min(start + max_per_grid, len(image_files))
        chunk = image_files[start:end]

        # Build grid for this chunk
        grid_image = assemble_grid(
            chunk, columns, width, start, text_regions, slide_size
        )

        # Determine output filename
        if len(image_files) <= max_per_grid:
            # Single grid - no suffix
            grid_file = target_path
        else:
            # Multiple grids - add index
            base = target_path.stem
            ext = target_path.suffix
            grid_file = target_path.parent / f"{base}-{chunk_idx + 1}{ext}"

        # Save grid
        grid_file.parent.mkdir(parents=True, exist_ok=True)
        grid_image.save(str(grid_file), quality=IMAGE_QUALITY)
        output_files.append(str(grid_file))

    return output_files


def assemble_grid(
    image_files,
    columns,
    width,
    start_slide=0,
    text_regions=None,
    slide_size=None,
):
    """Assemble preview grid from slide images with optional text region highlighting."""
    text_size = int(width * TEXT_SCALE)
    text_margin = int(text_size * TEXT_MARGIN_SCALE)

    # Calculate dimensions
    with Image.open(image_files[0]) as img:
        ratio = img.height / img.width
    height = int(width * ratio)

    # Determine grid size
    rows = (len(image_files) + columns - 1) // columns
    grid_width = columns * width + (columns + 1) * SPACING
    grid_height = rows * (height + text_size + text_margin * 2) + (rows + 1) * SPACING

    # Create canvas
    canvas = Image.new("RGB", (grid_width, grid_height), "white")
    draw = ImageDraw.Draw(canvas)

    # Load font
    try:
        font = ImageFont.load_default(size=text_size)
    except Exception:
        font = ImageFont.load_default()

    # Position previews
    for i, img_file in enumerate(image_files):
        row, col = i // columns, i % columns
        x = col * width + (col + 1) * SPACING
        y_start = row * (height + text_size + text_margin * 2) + (row + 1) * SPACING

        # Add slide number label
        label_text = f"{start_slide + i}"
        bbox = draw.textbbox((0, 0), label_text, font=font)
        text_w = bbox[2] - bbox[0]
        draw.text(
            (x + (width - text_w) // 2, y_start + text_margin),
            label_text,
            fill="black",
            font=font,
        )

        # Position preview below label
        y_preview = y_start + text_margin + text_size + text_margin

        with Image.open(img_file) as img:
            orig_w, orig_h = img.size

            # Apply text region outlines if enabled
            if text_regions and (start_slide + i) in text_regions:
                if img.mode != "RGBA":
                    img = img.convert("RGBA")

                regions = text_regions[start_slide + i]

                # Calculate scale factors
                if slide_size:
                    slide_w_inches, slide_h_inches = slide_size
                else:
                    slide_w_inches = orig_w / EXPORT_DPI
                    slide_h_inches = orig_h / EXPORT_DPI

                x_scale = orig_w / slide_w_inches
                y_scale = orig_h / slide_h_inches

                # Create highlight layer
                highlight = Image.new("RGBA", img.size, (255, 255, 255, 0))
                highlight_draw = ImageDraw.Draw(highlight)

                # Draw red borders around text regions
                for region in regions:
                    px_left = int(region["left"] * x_scale)
                    px_top = int(region["top"] * y_scale)
                    px_width = int(region["width"] * x_scale)
                    px_height = int(region["height"] * y_scale)

                    stroke = max(5, min(orig_w, orig_h) // 150)
                    highlight_draw.rectangle(
                        [(px_left, px_top), (px_left + px_width, px_top + px_height)],
                        outline=(255, 0, 0, 255),
                        width=stroke,
                    )

                img = Image.alpha_composite(img, highlight)
                img = img.convert("RGB")

            img.thumbnail((width, height), Image.Resampling.LANCZOS)
            w, h = img.size
            offset_x = x + (width - w) // 2
            offset_y = y_preview + (height - h) // 2
            canvas.paste(img, (offset_x, offset_y))

            # Add border
            if BORDER_THICKNESS > 0:
                draw.rectangle(
                    [
                        (offset_x - BORDER_THICKNESS, offset_y - BORDER_THICKNESS),
                        (
                            offset_x + w + BORDER_THICKNESS - 1,
                            offset_y + h + BORDER_THICKNESS - 1,
                        ),
                    ],
                    outline="gray",
                    width=BORDER_THICKNESS,
                )

    return canvas


if __name__ == "__main__":
    run()
